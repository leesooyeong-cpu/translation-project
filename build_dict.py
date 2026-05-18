import sys
import io
import re
import json
import time
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# ── 용어 규칙 (최우선) ────────────────────────────────────────────────
TERM_RULES = {
    "Dập Tự Động": "압착기",
    "dập tự động": "압착기",
    "Dập tự động": "압착기",
    "DAP TU DONG": "압착기",
}

# ── 베트남어 문자 판별 ────────────────────────────────────────────────
VN_CHARS = set(
    'àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩ'
    'òóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđ'
    'ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨ'
    'ÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ'
    'ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨ'
)

def is_vietnamese(text):
    return bool(text.strip()) and any(c in VN_CHARS for c in text)

def extract_texts_from_shape(shape, collected):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = ''.join(r.text for r in para.runs).strip()
            if line:
                collected.add(line)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    line = ''.join(r.text for r in para.runs).strip()
                    if line:
                        collected.add(line)
    if shape.shape_type == 6:  # GROUP
        for s in shape.shapes:
            extract_texts_from_shape(s, collected)

def extract_from_pptx(path):
    collected = set()
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                extract_texts_from_shape(shape, collected)
    except Exception as e:
        print(f"  [추출 오류] {Path(path).name}: {e}")
    return collected

def tokenize_vn_words(texts):
    """문장에서 베트남어 단어/구 추출"""
    words = set()
    for text in texts:
        if not is_vietnamese(text):
            continue
        # 구(phrase) 단위로 추가
        words.add(text.strip())
        # 개별 단어도 추가 (특수문자 기준 분리)
        parts = re.split(r'[/\\\(\)\[\]\{\}:;,\.\!\?\d]+', text)
        for part in parts:
            part = part.strip()
            if part and is_vietnamese(part) and len(part) > 1:
                words.add(part)
    return words

def batch_translate(word_list, batch_size=80):
    """여러 단어를 줄바꿈으로 묶어 한 번에 번역"""
    translator = GoogleTranslator(source='vi', target='ko')
    result = {}
    total = len(word_list)
    done = 0

    for i in range(0, total, batch_size):
        batch = word_list[i:i + batch_size]
        joined = '\n'.join(batch)
        try:
            translated = translator.translate(joined)
            lines = translated.split('\n')
            # 줄 수가 맞지 않으면 개별 번역으로 fallback
            if len(lines) != len(batch):
                for w in batch:
                    try:
                        t = translator.translate(w)
                        result[w] = t if t else w
                        time.sleep(0.05)
                    except Exception:
                        result[w] = w
            else:
                for w, t in zip(batch, lines):
                    result[w] = t.strip() if t.strip() else w
            time.sleep(0.3)
        except Exception as e:
            print(f"  [배치 오류] 개별 번역으로 전환: {e}")
            for w in batch:
                try:
                    t = translator.translate(w)
                    result[w] = t if t else w
                    time.sleep(0.1)
                except Exception:
                    result[w] = w

        done += len(batch)
        print(f"  번역 진행: {done}/{total}")

    return result

def main():
    src_root = Path(r"C:\hex_edit_project\VN_original")
    out_path = Path(r"C:\hex_edit_project\translation_dict.json")

    # 1. 파일 목록
    files = sorted(src_root.rglob("*.pptx")) + sorted(src_root.rglob("*.ppt"))
    total_files = len(files)
    print(f"총 {total_files}개 파일에서 텍스트 추출 중...\n")

    # 2. 텍스트 추출
    all_texts = set()
    for idx, f in enumerate(files, 1):
        texts = extract_from_pptx(str(f))
        all_texts.update(texts)
        if idx % 20 == 0 or idx == total_files:
            print(f"  [{idx}/{total_files}] 추출 완료, 누적 텍스트: {len(all_texts)}개")

    # 3. 베트남어 단어/구 추출 및 필터링
    print(f"\n단어 토크나이징 중...")
    vn_words = tokenize_vn_words(all_texts)
    print(f"고유 베트남어 항목: {len(vn_words)}개")

    # 4. 용어 규칙 먼저 적용
    translation_dict = {}
    for term, ko in TERM_RULES.items():
        translation_dict[term] = ko

    # 5. 나머지 번역 대상 추출
    to_translate = [w for w in sorted(vn_words) if w not in translation_dict]
    print(f"번역 대상: {len(to_translate)}개 (용어 규칙 {len(TERM_RULES)}개 제외)\n")

    # 6. 배치 번역
    print("Google 번역 시작...")
    translated = batch_translate(to_translate)
    translation_dict.update(translated)

    # 7. 용어 규칙을 최우선으로 덮어쓰기 (혹시 batch에서 덮였을 경우 대비)
    for term, ko in TERM_RULES.items():
        translation_dict[term] = ko

    # 8. 저장 (키 기준 정렬)
    sorted_dict = dict(sorted(translation_dict.items()))
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(sorted_dict, f, ensure_ascii=False, indent=2)

    print(f"\n완료! translation_dict.json 저장: {len(sorted_dict)}개 항목")
    print(f"경로: {out_path}")

if __name__ == "__main__":
    main()
