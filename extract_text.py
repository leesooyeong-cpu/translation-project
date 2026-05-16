import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

def extract_text(path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        print(f"=== Slide {i+1} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        print(f"  [{shape.shape_type}] {text}")
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            print(f"  [TABLE] {text}")

extract_text(sys.argv[1])
