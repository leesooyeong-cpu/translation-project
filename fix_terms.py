import sys
import io
from pathlib import Path
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

REPLACEMENTS = {
    "스템핑 기계": "압착기",
}

def replace_in_tf(tf):
    changed = False
    for para in tf.paragraphs:
        for run in para.runs:
            for old, new in REPLACEMENTS.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    changed = True
    return changed

def fix_file(path):
    prs = Presentation(path)
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if replace_in_tf(shape.text_frame):
                    changed = True
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if replace_in_tf(cell.text_frame):
                            changed = True
            if shape.shape_type == 6:
                for s in shape.shapes:
                    if s.has_text_frame:
                        if replace_in_tf(s.text_frame):
                            changed = True
    if changed:
        prs.save(path)
    return changed

def main():
    out_dir = Path(r"C:\hex_edit_project\KO_result")
    files = sorted(out_dir.glob("*_ko.pptx"))
    total = len(files)
    updated = 0

    for idx, f in enumerate(files, 1):
        try:
            changed = fix_file(str(f))
            if changed:
                updated += 1
                print(f"[{idx}/{total}] 수정: {f.name}")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {f.name} → {e}")

    print(f"\n완료. {updated}/{total}개 파일 수정됨.")

if __name__ == "__main__":
    main()
