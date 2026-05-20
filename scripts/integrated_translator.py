"""
베트남어 PPTX 번역기 (통합 버전)

사용법:
  python scripts/integrated_translator.py <소스_폴더_경로>
  python scripts/integrated_translator.py <소스_폴더_경로> --out <출력_폴더_경로>
  python scripts/integrated_translator.py <소스_폴더_경로> --dict <사전_파일_경로>

예시:
  python scripts/integrated_translator.py "C:/hex_edit_project/VN_original/045_Tiêu chuẩn_Công đoạn_Xoắn dây_done"
  python scripts/integrated_translator.py "C:/hex_edit_project/VN_original/080_Tiêu chuẩn_Công đoạn_Co nhiệt/TSWS-G080-RG3 CONT THETA3-00001"
"""

import sys
import io
import re
import json
import time
import argparse
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

DEFAULT_OUT_DIR  = Path(r"C:\hex_edit_project\KO_result")
DEFAULT_DICT_PATH = Path(r"C:\hex_edit_project\translation_dict.json")

# ─── 번역 상수 ──────────────────────────────────────────────────────────────

VN_CHARS = set(
    'àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩ'
    'òóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđ'
    'ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨ'
    'ÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ'
)

TERM_RULES = {
    "Dập Tự Động": "압착기",
    "dập tự động": "압착기",
    "DẬP TỰ ĐỘNG": "압착기",
    "Dập tự động": "압착기",
}

ACTION_VERBS = [
    '확인', '수행', '준비', '진행', '관리', '적용', '설정', '실시',
    '사용', '기록', '교체', '점검', '조정', '검사', '실행', '처리',
    '보관', '표시', '부착', '제거', '측정', '삽입', '연결', '분리',
    '청소', '유지', '보전', '작성', '등록', '배치', '정리', '고정',
    '투입', '장착', '탈착', '조립', '해체', '세척', '건조', '포장',
]

# ─── 파일명 매핑 (전체 누적) ─────────────────────────────────────────────────

FILENAME_MAP = {
    # ── 025 자동 절단 C01 ──
    "TSCS-G025-0001_A0_20220831_CẮT TỰ ĐỘNG_Danh sách kiểm tra hàng ngày ( sản xuất ).pptx":
        "TSCS-G025-0001_A0_20220831_자동 절단_일일 검사표 (생산)_ko.pptx",
    "TSIS-G025-0002_Tiêu chuẩn kiểm tra vỏ bọc cách điện của dây dẫn_a1.pptx":
        "TSIS-G025-0002_전선 절연 피복 검사 표준_a1_ko.pptx",
    "TSMS-G025-0003_Thông số trên màn hình điều khiển_a1.pptx":
        "TSMS-G025-0003_제어 화면 파라미터_a1_ko.pptx",
    "TSWS-G025-0003_Tiêu chuẩn xử lý (Khi máy xảy ra sự cố) (Cắt & uốn)_a1.pptx":
        "TSWS-G025-0003_처리 표준 (기계 고장 발생 시) (절단 및 굽힘)_a1_ko.pptx",
    # ── 025 자동 절단 C02 ──
    "TSCS-G025-0003_A0_20220831_C02_Danh sách kiểm tra hàng ngày ( sản xuất )( máy tời ).pptx":
        "TSCS-G025-0003_A0_20220831_C02_일일 검사표 (생산) (권취기)_ko.pptx",
    "TSMS-G025-0002_Tiêu Chuẩn Treo Linh Kiện Cắt – Dập_a0_thay hình ảnh.pptx":
        "TSMS-G025-0002_절단 압착 부품 걸이 표준_a0_ko.pptx",
    "TSPO-G025-0001_A1_20220831_Quy trình tuốt vỏ_C02.pptx":
        "TSPO-G025-0001_A1_20220831_피복 탈피 공정_C02_ko.pptx",
    # ── 026 차폐선 ──
    "TSCS-G026-0001_A1_20240610_Bảng Kiểm Tra Định Kỳ (Bảo Trì)_xoắn chống nhiễu.pptx":
        "TSCS-G026-0001_A1_20240610_정기 검사표 (보전)_차폐 꼬임선_ko.pptx",
    "TSCS-G026-0002-Phiếu kiểm tra 3D5S_a0 DÂY SHEILD.pptx":
        "TSCS-G026-0002_3D5S 검사표_차폐선_a0_ko.pptx",
    "TSCS-G026-0003_A1_20250218_DÂY SHIELD_Danh sách kiểm tra hàng ngày ( sản xuất ).pptx":
        "TSCS-G026-0003_A1_20250218_차폐선_일일 검사표 (생산)_ko.pptx",
    "TSCS-G026-0004_A0_20260505_Bảng Kiểm Tra Hằng Ngày (Sản Xuất).pptx":
        "TSCS-G026-0004_A0_20260505_일일 검사표 (생산)_ko.pptx",
    "TSCS-G026-0005_A0_20260505_Phiếu Kiểm Tra Định Kỳ (Bảo Trì).pptx":
        "TSCS-G026-0005_A0_20260505_정기 검사표 (보전)_ko.pptx",
    "TSIS-G026-0003_A0_20260505_Tiêu Chuẩn Bảo Vệ Đầu Terminal .pptx":
        "TSIS-G026-0003_A0_20260505_단자 헤드 보호 표준_ko.pptx",
    "TSMS-G026-0001_A0_20260505_Tiêu Chuẩn Nhận Dạng Phiếu Sản Xuất.pptx":
        "TSMS-G026-0001_A0_20260505_생산표 식별 표준_ko.pptx",
    "TSMS-G026-0002_A0_20260505_Tiêu Chuẩn Nhận Dạng Vật Tư.pptx":
        "TSMS-G026-0002_A0_20260505_자재 식별 표준_ko.pptx",
    "TSMS-G026-0003_A0_20260505_Tiêu Chuẩn Treo Linh Kiện Cắt – Dập.pptx":
        "TSMS-G026-0003_A0_20260505_절단 압착 부품 걸이 표준_ko.pptx",
    "TSNC-G026-0001_A0_20260505_Quy Trình Xử Lý Sản Phẩm Không Phù Hợp.pptx":
        "TSNC-G026-0001_A0_20260505_부적합 제품 처리 절차_ko.pptx",
    "TSPO-G026-0002_A0_20260505_Quy trình làm việc cắt dây sheild.pptx":
        "TSPO-G026-0002_A0_20260505_차폐선 절단 작업 공정_ko.pptx",
    "TSWS-G026-00010_A0_20260505_Tiêu Chuẩn Khắc Phục Khi Máy Xảy Ra Vấn Đề.pptx":
        "TSWS-G026-00010_A0_20260505_설비 문제 발생 시 복구 표준_ko.pptx",
    "TSWS-G026-00011_A0_20260505_Tiêu Chuẩn Vệ Sinh (3C5S).pptx":
        "TSWS-G026-00011_A0_20260505_위생 표준 (3C5S)_ko.pptx",
    "TSWS-G026-0007_A0_20260505_Tiêu Chuẩn Sử Dụng Máy Cắt Dây Shield.pptx":
        "TSWS-G026-0007_A0_20260505_차폐선 절단기 사용 표준_ko.pptx",
    "TSWS-G026-0008_A0_20260505_Tiêu Chuẩn Cắt Dây Shield.pptx":
        "TSWS-G026-0008_A0_20260505_차폐선 절단 표준_ko.pptx",
    "TSWS-G026-0009_A0_20260505_Tiêu Chuẩn Sử Dụng Phiếu Thông Tin Linh Kiện Lỗi.pptx":
        "TSWS-G026-0009_A0_20260505_불량 부품 정보표 사용 표준_ko.pptx",
    # ── 027 최종 피복 탈피 ──
    "TSCS-G027-0001_A1_20240610_TUỐT VỎ CUỐI_Bảng kiểm tra hàng ngày  (sản xuất )(Xưởng ).pptx":
        "TSCS-G027-0001_A1_20240610_최종 피복 탈피_일일 검사표 (생산) (현장)_ko.pptx",
    "TSCS-G027-0003_A2_20250408_Bảng kiểm tra hàng ngày( ABAG + Rework).pptx":
        "TSCS-G027-0003_A2_20250408_일일 검사표 (ABAG + 재작업)_ko.pptx",
    # ── 045 전선 꼬임 ──
    "TC máy xoắn dây.pptx":
        "TC_전선 꼬임 기계_ko.pptx",
    "thước đo dây xoắn đoạn ko xoắn bản mới.pptx":
        "비꼬임 구간 꼬임 전선 측정 도구 (신규)_ko.pptx",
    "TSCS-G045-0001_A2_20251029.XOẮN DÂY_Phiếu kiểm tra hàng ngày (Sản xuất).pptx":
        "TSCS-G045-0001_A2_20251029_전선 꼬임_일일 검사표 (생산)_ko.pptx",
    "TSCS-G045-0002_A0_20240610_Phiếu kiểm tra định kỳ ( Bảo trì).pptx":
        "TSCS-G045-0002_A0_20240610_정기 검사표 (보전)_ko.pptx",
    "TSCS-G045-0003_A1_20240610_Phiếu kiểm tra hàng ngày_Delta-.pptx":
        "TSCS-G045-0003_A1_20240610_일일 검사표_Delta-_ko.pptx",
    "TSCS-G045-0003_A2_20251029_Phiếu kiểm tra hàng ngày_Delta.pptx":
        "TSCS-G045-0003_A2_20251029_일일 검사표_Delta_ko.pptx",
    "TSIS-G045-000001_A1_20240610_Tiêu chuẩn tự kiểm tra dây xoắn.pptx":
        "TSIS-G045-000001_A1_20240610_전선 꼬임 자주 검사 표준_ko.pptx",
    "TSIS-G045-0001_A1_20240610_Tiêu chuẩn tự kiểm tra dây xoắn.pptx":
        "TSIS-G045-0001_A1_20240610_전선 꼬임 자주 검사 표준_ko.pptx",
    "TSMS-G045-0001_A2_20240610_Tiêu Chuẩn Quản Lý Xoắn DâyFLALUY.pptx":
        "TSMS-G045-0001_A2_20240610_전선 꼬임 관리 표준_FLALUУ_ko.pptx",
    "TSMS-G045-0001_A3_20250708_Tiêu Chuẩn Quản Lý Xoắn DâyFLALUY.pptx":
        "TSMS-G045-0001_A3_20250708_전선 꼬임 관리 표준_FLALUУ_ko.pptx",
    "TSMS-G045-0002_A2_20250816_Tiêu Chuẩn Cốc Bảo Vệ Đầu Terminal.pptx":
        "TSMS-G045-0002_A2_20250816_단자 헤드 보호컵 표준_ko.pptx",
    "TSNC-G045-0001_A1_20240610_Quy Trình Xử Lý Sản Phẩm Không Phù Hợp.pptx":
        "TSNC-G045-0001_A1_20240610_부적합 제품 처리 절차_ko.pptx",
    "TSPO-G045-0001_A2_2020708_Quy trình làm việc xoắn.pptx":
        "TSPO-G045-0001_A2_2020708_꼬임 작업 공정_ko.pptx",
    "TSWS-G045-0001_A1_20240610_Tiêu chuẩn sử dụng bảng điều khiển xoắn dây.pptx":
        "TSWS-G045-0001_A1_20240610_전선 꼬임 제어판 사용 표준_ko.pptx",
    "TSWS-G045-0001_A1_20240610_Tiêu chuẩn sử dụng bảng điều khiển xoắn dây_DELTA.pptx":
        "TSWS-G045-0001_A1_20240610_전선 꼬임 제어판 사용 표준_DELTA_ko.pptx",
    "TSWS-G045-0002_A2_20240610_Tiêu chuẩn độ dài xoắn dây (2).pptx":
        "TSWS-G045-0002_A2_20240610_전선 꼬임 길이 표준 (2)_ko.pptx",
    "TSWS-G045-0002_A2_20240610_Tiêu chuẩn độ dài xoắn dây (2)_DELTA.pptx":
        "TSWS-G045-0002_A2_20240610_전선 꼬임 길이 표준 (2)_DELTA_ko.pptx",
    "TSWS-G045-0003_A1_20240610_Tiêu chuẩn điều chỉnh vị trí phần căng thẳng dây.pptx":
        "TSWS-G045-0003_A1_20240610_전선 장력부 위치 조정 표준_ko.pptx",
    "TSWS-G045-0003_A1_20240610_Tiêu chuẩn điều chỉnh vị trí phần căng thẳng dây_DELTA.pptx":
        "TSWS-G045-0003_A1_20240610_전선 장력부 위치 조정 표준_DELTA_ko.pptx",
    "TSWS-G045-0004_Tiêu chuẩn xử lý (Khi máy xảy ra sự cố) _a1.pptx":
        "TSWS-G045-0004_처리 표준 (기계 고장 발생 시)_a1_ko.pptx",
    "TSWS-G045-0004_Tiêu chuẩn xử lý (Khi máy xảy ra sự cố) _a1_DELTA.pptx":
        "TSWS-G045-0004_처리 표준 (기계 고장 발생 시)_a1_DELTA_ko.pptx",
    # ── 050 중간 피복 탈피 ──
    "TSCS-G050-0001_A2_20250210_Phiếu kiểm tra hàng ngày (sản xuất).ppt":
        "TSCS-G050-0001_A2_20250210_일일 검사표 (생산)_ko.ppt",
    "TSCS-G050-0002_A2_20250210_Phiếu kiểm tra định kỳ (bảo trì).pptx":
        "TSCS-G050-0002_A2_20250210_정기 검사표 (보전)_ko.pptx",
    "TSIO-G050-0001_A2_20250210_Quy trình kiểm tra.pptx":
        "TSIO-G050-0001_A2_20250210_검사 절차_ko.pptx",
    "TSIS-G050-0001_A2_20250210_Tiêu Chuẩn Kiểm Tra Công Đoạn Tuốt Vỏ Giữa.pptx":
        "TSIS-G050-0001_A2_20250210_중간 피복 탈피 공정 검사 표준_ko.pptx",
    "TSPO-G050-0001_A3_20250210_Quy trình tuốt vỏ giữa.pptx":
        "TSPO-G050-0001_A3_20250210_중간 피복 탈피 공정_ko.pptx",
    "TSWS-G050-0001_A2_20250210_Tiêu Chuẩn Chuẩn Bị Làm Việc.pptx":
        "TSWS-G050-0001_A2_20250210_작업 준비 표준_ko.pptx",
    "TSWS-G050-0002_A2_20250210_Tiêu Chuẩn Thao Tác Bảo Vệ Cho Các Linh Kiện Cắt Dập.pptx":
        "TSWS-G050-0002_A2_20250210_절단 압착 부품 보호 작업 표준_ko.pptx",
    "TSWS-G050-0003_A2_20250210_Tiêu Chuẩn Thay Dao.pptx":
        "TSWS-G050-0003_A2_20250210_칼날 교체 표준_ko.pptx",
    "TSWS-G050-0004_A2_20250210_Tiêu Chuẩn Tuốt Vỏ Giữa.pptx":
        "TSWS-G050-0004_A2_20250210_중간 피복 탈피 표준_ko.pptx",
    "TSWS-G050-0005_A2_20250210_Phương pháp quấn miếng bảo vệ linh kiện dập.pptx":
        "TSWS-G050-0005_A2_20250210_압착 부품 보호 테이프 감기 방법_ko.pptx",
    # ── 065 초음파 용착 ──
    "TSCS-G065-0001_A1_20240610_HÀN SIÊU ÂM_Phiếu kiểm tra ( sản xuất ).pptx":
        "TSCS-G065-0001_A1_20240610_초음파 용착_일일 검사표 (생산)_ko.pptx",
    "TSCS-G065-0002_A1_20240610_Danh sách kiểm tra thường xuyên (Bảo trì).pptx":
        "TSCS-G065-0002_A1_20240610_정기 검사표 (보전)_ko.pptx",
    "TSIO-G065-0001_A1_20240610_Quá trình kiểm tra hàn siêu âm.pptx":
        "TSIO-G065-0001_A1_20240610_초음파 용착 검사 절차_ko.pptx",
    "TSIS-G065-0001_A9_20251124_Tiêu chuẩn đo CCH mối hàn siêu âm_new2025.pptx":
        "TSIS-G065-0001_A9_20251124_초음파 용착부 CCH 측정 표준_new2025_ko.pptx",
    "TSIS-G065-0002_A1_20240610_Tiêu chuẩn tự kiểm tra của hàn siêu âm.pptx":
        "TSIS-G065-0002_A1_20240610_초음파 용착 자주 검사 표준_ko.pptx",
    "TSMS-G065-0001_A1_20240610_Tiêu Chuẩn Cốc Bảo Vệ Đầu Terminal.pptx":
        "TSMS-G065-0001_A1_20240610_단자 헤드 보호컵 표준_ko.pptx",
    "TSMS-G065-0002_Tiêu Chuẩn Treo Linh Kiện Cắt.pptx":
        "TSMS-G065-0002_절단 부품 걸이 표준_ko.pptx",
    "TSNC-G065-0001_a1_quy trình xử lý linh kiện không phù hợp.pptx":
        "TSNC-G065-0001_a1_부적합 부품 처리 절차_ko.pptx",
    "TSPO-G065-0001_A1_20240610_Quy trình hàn siêu âm.pptx":
        "TSPO-G065-0001_A1_20240610_초음파 용착 공정_ko.pptx",
    "TSWS-G065-0001_A1_20240610_Phương pháp quấn miếng bảo vệ linh kiện dập.pptx":
        "TSWS-G065-0001_A1_20240610_압착 부품 보호 테이프 감기 방법_ko.pptx",
    "TSWS-G065-0002_Tiêu chuẩn sử dụng thước đo điện tử (micromet)_a2.pptx":
        "TSWS-G065-0002_전자 측정기(마이크로미터) 사용 표준_a2_ko.pptx",
    # ── 070 테이프핑 ──
    "TSCS-G070-0001_A1_20240610_QUẤN KEO_Phiếu kiểm tra (sản xuất).pptx":
        "TSCS-G070-0001_A1_20240610_테이프핑_일일 검사표 (생산)_ko.pptx",
    "TSCS-G070-0002_A1_20240610_Danh sách kiểm tra định kỳ (Bảo trì).pptx":
        "TSCS-G070-0002_A1_20240610_정기 검사표 (보전)_ko.pptx",
    "TSMS-G070-0001_A0_20220831_Tiêu chuẩn quản lý băng keo.pptx":
        "TSMS-G070-0001_A0_20220831_접착 테이프 관리 표준_ko.pptx",
    "TSPO-G070-0001_A1_20240610_Quy trình quấn keo.pptx":
        "TSPO-G070-0001_A1_20240610_테이프 감기 공정_ko.pptx",
    "TSWS-G070-0001_A1_20240610_Quy trình khai thác của khớp nối.pptx":
        "TSWS-G070-0001_A1_20240610_조인트 작업 절차_ko.pptx",
    "TSWS-G070-0002_A1_20240610_Tiêu chuẩn quy trình quấn keo.pptx":
        "TSWS-G070-0002_A1_20240610_테이프 감기 공정 표준_ko.pptx",
    "TSWS-G070-0003_A0 Tiêu chuẩn quy trình thay cuộn keo_a0.pptx":
        "TSWS-G070-0003_A0_접착 테이프 롤 교체 공정 표준_a0_ko.pptx",
    "TSWS-G070-0004_A1_20250607_Tiêu chuẩn thao tác máy quấn băng dính.pptx":
        "TSWS-G070-0004_A1_20250607_테이프 감기 기계 조작 표준_ko.pptx",
    # ── 080 열수축 (메인 폴더) ──
    "TSCS-G080-0001_A1_20240610_CO NHIỆT_Phiếu kiểm tra ( sản xuất )_H-01.ppt":
        "TSCS-G080-0001_A1_20240610_열수축_검사표 (생산)_H-01_ko.pptx",
    "TSCS-G080-0002_A1_20240610_CO NHIỆT_Phiếu kiểm tra ( sản xuất )_H-02.ppt":
        "TSCS-G080-0002_A1_20240610_열수축_검사표 (생산)_H-02_ko.pptx",
    "TSCS-G080-0003_A1_20240610_Phiếu kiểm tra ( bảo trì ).ppt":
        "TSCS-G080-0003_A1_20240610_검사표 (보전)_ko.pptx",
    "TSIS-G080-0001_a2_20240822_Tiêu chuẩn tự kiểm tra khớp nối.pptx":
        "TSIS-G080-0001_a2_20240822_조인트 자가검사 표준_ko.pptx",
    "TSIS-G080-0002_a2_Tiêu chuẩn tự kiểm tra Ring terminal.pptx":
        "TSIS-G080-0002_a2_링 터미널 자가검사 표준_ko.pptx",
    "TSIS-G080-0003_Tiêu chuẩn tự kiểm tra của ống co nhiệt (sóng siêu âm)_a1.pptx":
        "TSIS-G080-0003_열수축 튜브 자가검사 표준 (초음파)_a1_ko.pptx",
    "TSMS_G080-0001_A0_20250217_Tiêu chuẩn quản lý ống co nhiệt_(+dán bàn co nhiệt).pptx":
        "TSMS_G080-0001_A0_20250217_열수축 튜브 관리 표준_(+게시판 부착)_ko.pptx",
    "TSMS-G080-0002_A1_20240610_Tiêu chuẩn nhận dạng phiếu sản xuất.pptx":
        "TSMS-G080-0002_A1_20240610_생산표 식별 표준_ko.pptx",
    "TSMS-G080-0003_A1_20240610_Tiêu chuẫn nhiệt độ và tốc độ máy co nhiệt HT-01+02(+dán bàn).pptx":
        "TSMS-G080-0003_A1_20240610_열수축기 HT-01+02 온도 및 속도 표준_(+게시판 부착)_ko.pptx",
    "TSMS-G080-0004_A1_20240610_Tiêu chuẫn nhiệt độ que đo nhiệt(DÁN BÀN ).pptx":
        "TSMS-G080-0004_A1_20240610_온도 프로브 온도 표준_(게시판 부착)_ko.pptx",
    "TSMS-G080-0005_Tiêu chuẩn phân biệt ống co nhiệt_a2.pptx":
        "TSMS-G080-0005_열수축 튜브 구분 표준_a2_ko.pptx",
    "TSMS-G080-0006_Tiêu chuẩn cốc bảo vệ_a3.pptx":
        "TSMS-G080-0006_보호 캡 표준_a3_ko.pptx",
    "TSMS-G080-0007_Tiêu chuẩn quản lý co nhiệt_a3.pptx":
        "TSMS-G080-0007_열수축 관리 표준_a3_ko.pptx",
    "TSMS-G080-0008_Tiêu chuẩn giữ vật phẩm cắt-dập (Khi cất giữ trong hộp)_a0.pptx":
        "TSMS-G080-0008_절단-압착 부품 보관 표준 (박스 보관 시)_a0_ko.pptx",
    "TSMS-G080-0009_A2_Tiêu Chuẩn Treo Linh Kiện Cắt – Dập.pptx":
        "TSMS-G080-0009_A2_절단-압착 부품 걸기 표준_ko.pptx",
    "TSMS-G080-0010_Tiêu chuẩn quản lý đặc tính đặc biệt (nhiệt độ  tốc độ)_a2.pptx":
        "TSMS-G080-0010_특수 특성 관리 표준 (온도 속도)_a2_ko.pptx",
    "TSNC-G080-0001_A1_20240610_Quy Trình Xử Lý Sản Phẩm Không Phù Hợp.pptx":
        "TSNC-G080-0001_A1_20240610_부적합 제품 처리 절차_ko.pptx",
    "TSPO-G080-0001_a3_Quy trình làm việc co nhiệt .pptx":
        "TSPO-G080-0001_a3_열수축 작업 절차_ko.pptx",
    "TSPO-G080-0002_Quy trình kiểm tra co nhiệt_a2.pptx":
        "TSPO-G080-0002_열수축 검사 절차_a2_ko.pptx",
    "TSWM-G080-0009_HT02_A3_20240622_Tiêu chuẩn co nhiệt (HT-02+01_TK1).pptx":
        "TSWM-G080-0009_HT02_A3_20240622_열수축 표준 (HT-02+01_TK1)_ko.pptx",
    "TSWM-G080-HT01_HT01_A5_20251110_Tiêu chuẩn co nhiệt (HT-01)(DÁN BÀN).pptx":
        "TSWM-G080-HT01_HT01_A5_20251110_열수축 표준 (HT-01)(게시판 부착)_ko.pptx",
    "TSWM-G080-HT02_HT02_A5_20251110_Tiêu chuẩn co nhiệt (HT-02)(DÁN BÀN).pptx":
        "TSWM-G080-HT02_HT02_A5_20251110_열수축 표준 (HT-02)(게시판 부착)_ko.pptx",
    "TSWS-G080_0001_A0_20240115_Phương pháp mở biểu đồ co nhiệt(DÁN BÀN).pptx":
        "TSWS-G080_0001_A0_20240115_열수축 차트 열기 방법_(게시판 부착)_ko.pptx",
    "TSWS-G080_0007_A0_20240115_Phương pháp scan vật tư co nhiệt.pptx":
        "TSWS-G080_0007_A0_20240115_열수축 자재 스캔 방법_ko.pptx",
    "TSWS-G080-0002_A0_20240610_Tiêu chuẩn sử dụng nhiệt kế.pptx":
        "TSWS-G080-0002_A0_20240610_온도계 사용 표준_ko.pptx",
    "TSWS-G080-0003_A1_20240610_Tiêu chuẩn quản lý của quá trình ép nhiệt (quản lý SPC)_làm lại theo TS.pptx":
        "TSWS-G080-0003_A1_20240610_열압착 공정 관리 표준 (SPC 관리)_ko.pptx",
    "TSWS-G080-0004_A1_20240610_Tiêu chuẩn chèn ống co nhiệt.pptx":
        "TSWS-G080-0004_A1_20240610_열수축 튜브 삽입 표준_ko.pptx",
    "TSWS-G080-0005_A2_20251110_Tiêu chuẩn mẫu thử ép nhiệt.pptx":
        "TSWS-G080-0005_A2_20251110_열압착 시험 샘플 표준_ko.pptx",
    "TSWS-G080-0006_A1_20240610_Tiêu Chuẩn Khắc Phục (Khi Máy Xảy Ra Vấn Đề).pptx":
        "TSWS-G080-0006_A1_20240610_복구 표준 (설비 문제 발생 시)_ko.pptx",
    "TSWS-G080-0008_A0_Tiêu chuẩn sử dụng nhiệt kế _máy khò nhiệt.pptx":
        "TSWS-G080-0008_A0_온도계 사용 표준_열풍기_ko.pptx",
    "TSWS-G080-RG3 CONT THETA3-00001_a1_TIÊU CHUẨN LÀM VIỆC SÚNG BẮN NHIỆT.pptx":
        "TSWS-G080-RG3_CONT_THETA3-00001_a1_열풍건 작업 표준_ko.pptx",
    # ── 080 RG3 CONT THETA3-00001 폴더 ──
    "TSCS-G080-0001_A1_20240610_Phiếu kiểm tra ( sản xuất )_khò nhiệt.ppt":
        "TSCS-G080-0001_A1_20240610_검사표 (생산)_열풍기_ko.pptx",
    "TSIO-G080-0001_A1_20240610_Quy trình kiểm tra co nhiệt.pptx":
        "TSIO-G080-0001_A1_20240610_열수축 검사 절차_ko.pptx",
    "TSMS-G080-0003_A1_20240610_Tiêu chuẫn nhiệt độ máy khò nhiệt.pptx":
        "TSMS-G080-0003_A1_20240610_열풍기 온도 표준_ko.pptx",
    "TSPO-G080-0001_A2_20240610_Quy trình làm việc co nhiệt dạng (RING  SPLICE JOINT ).pptx":
        "TSPO-G080-0001_A2_20240610_열수축 작업 절차_RING_SPLICE_JOINT_ko.pptx",
    "TSWS-G080-00011_a1Tiêu Chuẩn Thao Tác Bảo Vệ Cho Các Linh Kiện Cắt Dập.pptx":
        "TSWS-G080-00011_a1_절단 압착 부품 보호 작업 표준_ko.pptx",
    # ── 085 내전압 시험 ──
    "TSCS-G085-0001_Báo cáo sản phẩm lỗi hàng ngày.pptx":
        "TSCS-G085-0001_일일 불량품 보고서_ko.pptx",
    "TSPO-GS085-0001_A0_20220831_Quy trình kiểm tra điện áp chịu được.pptx":
        "TSPO-GS085-0001_A0_20220831_내전압 검사 공정_ko.pptx",
    "TSWS-G085-0001_A1_20240610_Bảo vệ tiêu chuẩn công việc của mặt hàng cắt & uốn.pptx":
        "TSWS-G085-0001_A1_20240610_절단 및 굽힘 부품 보호 작업 표준_ko.pptx",
    "TSWS-G085-0002_A1_20240610_Tiêu chuẩn kiểm tra điện áp chịu được (TOS5200).pptx":
        "TSWS-G085-0002_A1_20240610_내전압 검사 표준 (TOS5200)_ko.pptx",
    "TSWS-G085-0003_A1_20240610_Phương pháp thiết lập tình trạng máy kiểm tra điện áp chịu được.pptx":
        "TSWS-G085-0003_A1_20240610_내전압 시험기 조건 설정 방법_ko.pptx",
    # ── 100 집합 공정 ──
    "TSWS-G100-0002_Tiêu chuẩn treo mạch xe tập kết DN8aPE.pptx":
        "TSWS-G100-0002_집합 대차 배선 걸이 표준_DN8aPE_ko.pptx",
}


# ─── 핵심 번역 함수 ──────────────────────────────────────────────────────────

def load_translation_dict(dict_path: Path) -> dict:
    try:
        with open(dict_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {}


def apply_dict(text: str, trans_dict: dict) -> str:
    for vn, ko in trans_dict.items():
        if vn in text:
            text = text.replace(vn, ko)
    return text


def is_vietnamese(text: str) -> bool:
    return any(c in VN_CHARS for c in text)


def to_concise(text: str) -> str:
    t = text.strip()
    if len(t) < 3:
        return t
    for v in ACTION_VERBS:
        for ending in ['합니다', '한다', '하다']:
            t = re.sub(rf'([을를])\s*{v}{ending}\.?$', f' {v}.', t)
        t = re.sub(rf'([을를])\s*{v}해야\s*합니다\.?$', f' {v}해야 함.', t)
    endings = [
        (r'해야\s*합니다\.?$',   '해야 함.'),
        (r'해야\s*한다\.?$',     '해야 함.'),
        (r'하여야\s*합니다\.?$', '해야 함.'),
        (r'하십시오\.?$',        '할 것.'),
        (r'하세요\.?$',          '할 것.'),
        (r'해주세요\.?$',        '할 것.'),
        (r'합니다\.?$',          '함.'),
        (r'입니다\.?$',          '임.'),
        (r'됩니다\.?$',          '됨.'),
        (r'있습니다\.?$',        '있음.'),
        (r'없습니다\.?$',        '없음.'),
        (r'한다\.?$',            '함.'),
        (r'된다\.?$',            '됨.'),
        (r'있다\.?$',            '있음.'),
        (r'없다\.?$',            '없음.'),
        (r'이다\.?$',            '임.'),
    ]
    for pattern, replacement in endings:
        new_t = re.sub(pattern, replacement, t)
        if new_t != t:
            t = new_t
            break
    t = re.sub(r'([을를])\s*사용하여', ' 활용', t)
    t = re.sub(r'([을를])\s*이용하여', ' 활용', t)
    t = re.sub(r'\.{2,}$', '.', t)
    return t


def apply_term_rules(text: str) -> str:
    for vn, ko in TERM_RULES.items():
        text = text.replace(vn, ko)
    return text


def translate_text(text: str, translator, trans_dict: dict) -> str:
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    applied = apply_dict(text, trans_dict)
    if not is_vietnamese(applied):
        return apply_term_rules(to_concise(applied))
    try:
        result = translator.translate(applied)
        time.sleep(0.05)
        if result:
            result = apply_term_rules(result)
            result = to_concise(result)
        return result if result else text
    except Exception:
        return text


def process_text_frame(tf, translator, trans_dict: dict):
    for para in tf.paragraphs:
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator, trans_dict)
        elif len(para.runs) > 1:
            translated = translate_text(full_text, translator, trans_dict)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''


def process_shape(shape, translator, trans_dict: dict):
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, translator, trans_dict)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, translator, trans_dict)
    if shape.shape_type == 6:
        for s in shape.shapes:
            process_shape(s, translator, trans_dict)


def translate_pptx(input_path: str, output_path: str, trans_dict: dict):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, translator, trans_dict)
    prs.save(output_path)


# ─── 메인 ────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description='베트남어 PPTX → 한국어 번역 (통합 버전)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument('src', help='번역할 PPTX 파일이 있는 소스 폴더 경로')
    parser.add_argument('--out', default=str(DEFAULT_OUT_DIR),
                        help=f'번역 결과 저장 폴더 (기본값: {DEFAULT_OUT_DIR})')
    parser.add_argument('--dict', default=str(DEFAULT_DICT_PATH),
                        help=f'translation_dict.json 경로 (기본값: {DEFAULT_DICT_PATH})')
    args = parser.parse_args()

    src_dir = Path(args.src)
    out_dir = Path(args.out)
    dict_path = Path(args.dict)

    if not src_dir.exists():
        print(f"오류: 소스 폴더가 존재하지 않습니다 → {src_dir}")
        sys.exit(1)

    out_dir.mkdir(parents=True, exist_ok=True)

    trans_dict = load_translation_dict(dict_path)
    print(f"번역 사전 {len(trans_dict)}개 항목 로드 완료\n")

    files = sorted([
        f for f in src_dir.iterdir()
        if f.suffix.lower() in ('.pptx', '.ppt') and f.name != 'Thumbs.db'
    ])
    total = len(files)
    if total == 0:
        print("번역할 PPTX/PPT 파일이 없습니다.")
        sys.exit(0)
    print(f"총 {total}개 파일 번역 시작\n")

    errors, skipped = [], []

    for idx, f in enumerate(files, 1):
        # .ppt 구형 포맷 처리
        if f.suffix.lower() == '.ppt':
            print(f"[{idx}/{total}] 스킵 (.ppt 구형 포맷, python-pptx 미지원): {f.name}\n")
            skipped.append(f.name)
            continue

        # 출력 파일명 결정: 매핑 있으면 사용, 없으면 자동 생성
        out_name = FILENAME_MAP.get(f.name, f.stem + '_ko.pptx')
        if f.name not in FILENAME_MAP:
            print(f"[{idx}/{total}] 매핑 없음 → 자동 파일명: {out_name}")

        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name}\n        → {out_name} 번역 중...")
        try:
            translate_pptx(str(f), str(out_path), trans_dict)
            print(f"[{idx}/{total}] 완료\n")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {e}\n")
            errors.append(f"{f.name}: {e}")

    print("=" * 60)
    success = total - len(errors) - len(skipped)
    print(f"번역 완료. 성공: {success}개 / 오류: {len(errors)}개 / 스킵: {len(skipped)}개")
    if errors:
        print("오류 파일:")
        for e in errors:
            print(f"  - {e}")
    if skipped:
        print("스킵 파일 (.ppt 구형 포맷):")
        for s in skipped:
            print(f"  - {s}")


if __name__ == "__main__":
    main()
