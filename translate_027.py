import sys
import io
import re
import json
import time
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

VN_CHARS = set(
    'àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩ'
    'òóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđ'
    'ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨ'
    'ÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ'
)

TERM_RULES = {
    "Dập Tự Động": "압착기",
    "dập tự động": "압착기",
    "Dập tự động": "압착기",
}

ACTION_VERBS = [
    '확인', '수행', '준비', '진행', '관리', '적용', '설정', '실시',
    '사용', '기록', '교체', '점검', '조정', '검사', '실행', '처리',
    '보관', '표시', '부착', '제거', '측정', '삽입', '연결', '분리',
    '청소', '유지', '보전', '작성', '등록', '배치', '정리', '고정',
    '투입', '장착', '탈착', '조립', '해체', '세척', '건조', '포장',
]

FILENAME_MAP = {
    "TSCS-G027-0001_A1_20240610_TUỐT VỎ CUỐI_Bảng kiểm tra hàng ngày  (sản xuất )(Xưởng ).pptx":
        "TSCS-G027-0001_A1_20240610_최종 피복 탈피_일일 검사표 (생산) (현장)_ko.pptx",
    "TSCS-G027-0003_A2_20250408_Bảng kiểm tra hàng ngày( ABAG + Rework).pptx":
        "TSCS-G027-0003_A2_20250408_일일 검사표 (ABAG + 재작업)_ko.pptx",
}

def load_translation_dict(dict_path):
    try:
        with open(dict_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {}

def apply_dict(text, trans_dict):
    for vn, ko in trans_dict.items():
        if vn in text:
            text = text.replace(vn, ko)
    return text

def is_vietnamese(text):
    return any(c in VN_CHARS for c in text)

def to_concise(text: str) -> str:
    t = text.strip()
    if len(t) < 3:
        return t
    for v in ACTION_VERBS:
        for ending in ['합니다', '한다', '하다']:
            t = re.sub(rf'([을를])\s*{v}{ending}\.?$', f' {v}.', t)
        t = re.sub(rf'([을를])\s*{v}해야\s*합니다\.?$', f' {v}해야 함.', t)
    endings = [
        (r'해야\s*합니다\.?$',   '해야 함.'),
        (r'해야\s*한다\.?$',     '해야 함.'),
        (r'하여야\s*합니다\.?$', '해야 함.'),
        (r'하십시오\.?$',        '할 것.'),
        (r'하세요\.?$',          '할 것.'),
        (r'해주세요\.?$',        '할 것.'),
        (r'합니다\.?$',          '함.'),
        (r'입니다\.?$',          '임.'),
        (r'됩니다\.?$',          '됨.'),
        (r'있습니다\.?$',        '있음.'),
        (r'없습니다\.?$',        '없음.'),
        (r'한다\.?$',            '함.'),
        (r'된다\.?$',            '됨.'),
        (r'있다\.?$',            '있음.'),
        (r'없다\.?$',            '없음.'),
        (r'이다\.?$',            '임.'),
    ]
    for pattern, replacement in endings:
        new_t = re.sub(pattern, replacement, t)
        if new_t != t:
            t = new_t
            break
    t = re.sub(r'([을를])\s*사용하여', ' 활용', t)
    t = re.sub(r'([을를])\s*이용하여', ' 활용', t)
    t = re.sub(r'\.{2,}$', '.', t)
    return t

def apply_term_rules(text: str) -> str:
    for vn, ko in TERM_RULES.items():
        text = text.replace(vn, ko)
    return text

def translate_text(text, translator, trans_dict):
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    applied = apply_dict(text, trans_dict)
    if not is_vietnamese(applied):
        return apply_term_rules(to_concise(applied))
    try:
        result = translator.translate(applied)
        time.sleep(0.05)
        if result:
            result = apply_term_rules(result)
            result = to_concise(result)
        return result if result else text
    except Exception:
        return text

def process_text_frame(tf, translator, trans_dict):
    for para in tf.paragraphs:
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator, trans_dict)
        elif len(para.runs) > 1:
            translated = translate_text(full_text, translator, trans_dict)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''

def process_shape(shape, translator, trans_dict):
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, translator, trans_dict)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, translator, trans_dict)
    if shape.shape_type == 6:
        for s in shape.shapes:
            process_shape(s, translator, trans_dict)

def translate_pptx(input_path, output_path, trans_dict):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, translator, trans_dict)
    prs.save(output_path)

def main():
    src_dir = Path(r"C:\hex_edit_project\VN_original\027_Tiêu chuẩn_Công đoạn_Tuốt vỏ cuối")
    out_dir = Path(r"C:\hex_edit_project\KO_result")
    out_dir.mkdir(parents=True, exist_ok=True)

    dict_path = Path(r"C:\hex_edit_project\translation_dict.json")
    trans_dict = load_translation_dict(dict_path)
    print(f"번역 사전 {len(trans_dict)}개 항목 로드 완료\n")

    files = sorted([f for f in src_dir.iterdir() if f.suffix.lower() == '.pptx'])
    total = len(files)
    print(f"총 {total}개 파일 번역 시작\n")

    errors = []
    for idx, f in enumerate(files, 1):
        out_name = FILENAME_MAP.get(f.name)
        if not out_name:
            print(f"[{idx}/{total}] 파일명 매핑 없음, 스킵: {f.name}")
            errors.append(f.name)
            continue

        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name}\n        → {out_name} 번역 중...")
        try:
            translate_pptx(str(f), str(out_path), trans_dict)
            print(f"[{idx}/{total}] 완료\n")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {e}\n")
            errors.append(f.name)

    print("=" * 60)
    print(f"번역 완료. 성공: {total - len(errors)}개 / 실패: {len(errors)}개")
    if errors:
        for e in errors:
            print(f"  - {e}")
    return errors

if __name__ == "__main__":
    main()
