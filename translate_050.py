import os
import sys
import io
import time
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

VN_CHARS = set('àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩòóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ')

# 파일명 베트남어→한국어 매핑
FILENAME_MAP = {
    "TSCS-G050-0001_A2_20250210_Phiếu kiểm tra hàng ngày (sản xuất).ppt":
        "TSCS-G050-0001_A2_20250210_일일 검사표 (생산)_ko.ppt",
    "TSCS-G050-0002_A2_20250210_Phiếu kiểm tra định kỳ (bảo trì).pptx":
        "TSCS-G050-0002_A2_20250210_정기 검사표 (보전)_ko.pptx",
    "TSIO-G050-0001_A2_20250210_Quy trình kiểm tra.pptx":
        "TSIO-G050-0001_A2_20250210_검사 절차_ko.pptx",
    "TSIS-G050-0001_A2_20250210_Tiêu Chuẩn Kiểm Tra Công Đoạn Tuốt Vỏ Giữa.pptx":
        "TSIS-G050-0001_A2_20250210_중간 피복 탈피 공정 검사 표준_ko.pptx",
    "TSPO-G050-0001_A3_20250210_Quy trình tuốt vỏ giữa.pptx":
        "TSPO-G050-0001_A3_20250210_중간 피복 탈피 공정_ko.pptx",
    "TSWS-G050-0001_A2_20250210_Tiêu Chuẩn Chuẩn Bị Làm Việc.pptx":
        "TSWS-G050-0001_A2_20250210_작업 준비 표준_ko.pptx",
    "TSWS-G050-0002_A2_20250210_Tiêu Chuẩn Thao Tác Bảo Vệ Cho Các Linh Kiện Cắt Dập.pptx":
        "TSWS-G050-0002_A2_20250210_절단 압착 부품 보호 작업 표준_ko.pptx",
    "TSWS-G050-0003_A2_20250210_Tiêu Chuẩn Thay Dao.pptx":
        "TSWS-G050-0003_A2_20250210_칼날 교체 표준_ko.pptx",
    "TSWS-G050-0004_A2_20250210_Tiêu Chuẩn Tuốt Vỏ Giữa.pptx":
        "TSWS-G050-0004_A2_20250210_중간 피복 탈피 표준_ko.pptx",
    "TSWS-G050-0005_A2_20250210_Phương pháp quấn miếng bảo vệ linh kiện dập.pptx":
        "TSWS-G050-0005_A2_20250210_압착 부품 보호 테이프 감기 방법_ko.pptx",
}

def is_vietnamese(text):
    return any(c in VN_CHARS for c in text)

def translate_text(text, translator):
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    try:
        result = translator.translate(text)
        time.sleep(0.05)
        return result if result else text
    except Exception:
        return text

def process_text_frame(tf, translator):
    for para in tf.paragraphs:
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator)
        elif len(para.runs) > 1:
            translated = translate_text(full_text, translator)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''

def process_shape(shape, translator):
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, translator)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, translator)
    if shape.shape_type == 6:  # GROUP
        for s in shape.shapes:
            process_shape(s, translator)

def translate_pptx(input_path, output_path):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, translator)
    prs.save(output_path)

def main():
    src_dir = Path(r"C:\hex_edit_project\VN_original\050_Tiêu chuẩn_Công đoạn_Tuốt vỏ giữa")
    out_dir = Path(r"C:\베트남 티에스텍 업무\13. 절압팀\절차서 TC\050 표준_공정_중간 피복 탈피")
    out_dir.mkdir(parents=True, exist_ok=True)

    files = [f for f in src_dir.iterdir() if f.suffix.lower() in ('.pptx', '.ppt')]
    files.sort()
    total = len(files)
    print(f"총 {total}개 파일 번역 시작\n")

    skipped = []
    errors = []

    for idx, f in enumerate(files, 1):
        out_name = FILENAME_MAP.get(f.name)
        if not out_name:
            print(f"[{idx}/{total}] 파일명 매핑 없음, 스킵: {f.name}")
            skipped.append(f.name)
            continue

        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name}")
        print(f"        → {out_name} 번역 중...")

        if f.suffix.lower() == '.ppt':
            print(f"[{idx}/{total}] .ppt 형식은 python-pptx 미지원 → 파일 복사만 수행")
            import shutil
            shutil.copy2(str(f), str(out_path))
            print(f"[{idx}/{total}] 복사 완료 (내용 미번역): {out_name}")
            skipped.append(f.name + " (.ppt 형식 번역 불가, 복사만)")
            continue

        try:
            translate_pptx(str(f), str(out_path))
            print(f"[{idx}/{total}] 완료: {out_name}\n")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {f.name} → {e}\n")
            errors.append(f.name)

    print("=" * 60)
    print(f"전체 번역 완료. 성공: {total - len(skipped) - len(errors)}개")
    if skipped:
        print(f"예외/스킵: {len(skipped)}건")
        for s in skipped:
            print(f"  - {s}")
    if errors:
        print(f"오류: {len(errors)}건")
        for e in errors:
            print(f"  - {e}")

    return skipped, errors

if __name__ == "__main__":
    main()
