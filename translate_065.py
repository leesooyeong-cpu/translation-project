import os
import sys
import io
import re
import time
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

VN_CHARS = set(
    'àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩ'
    'òóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđ'
    'ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨ'
    'ÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ'
)

TERM_RULES = {
    "Dập Tự Động": "압착기",
    "dập tự động": "압착기",
    "Dập tự động": "압착기",
}

ACTION_VERBS = [
    '확인', '수행', '준비', '진행', '관리', '적용', '설정', '실시',
    '사용', '기록', '교체', '점검', '조정', '검사', '실행', '처리',
    '보관', '표시', '부착', '제거', '측정', '삽입', '연결', '분리',
    '청소', '유지', '보전', '작성', '등록', '배치', '정리', '고정',
    '투입', '장착', '탈착', '조립', '해체', '세척', '건조', '포장',
]

FILENAME_MAP = {
    "TSCS-G065-0001_A1_20240610_HÀN SIÊU ÂM_Phiếu kiểm tra ( sản xuất ).pptx":
        "TSCS-G065-0001_A1_20240610_초음파 용착_일일 검사표 (생산)_ko.pptx",
    "TSCS-G065-0002_A1_20240610_Danh sách kiểm tra thường xuyên (Bảo trì).pptx":
        "TSCS-G065-0002_A1_20240610_정기 검사표 (보전)_ko.pptx",
    "TSIO-G065-0001_A1_20240610_Quá trình kiểm tra hàn siêu âm.pptx":
        "TSIO-G065-0001_A1_20240610_초음파 용착 검사 절차_ko.pptx",
    "TSIS-G065-0001_A9_20251124_Tiêu chuẩn đo CCH mối hàn siêu âm_new2025.pptx":
        "TSIS-G065-0001_A9_20251124_초음파 용착부 CCH 측정 표준_new2025_ko.pptx",
    "TSIS-G065-0002_A1_20240610_Tiêu chuẩn tự kiểm tra của hàn siêu âm.pptx":
        "TSIS-G065-0002_A1_20240610_초음파 용착 자주 검사 표준_ko.pptx",
    "TSMS-G065-0001_A1_20240610_Tiêu Chuẩn Cốc Bảo Vệ Đầu Terminal.pptx":
        "TSMS-G065-0001_A1_20240610_단자 헤드 보호컵 표준_ko.pptx",
    "TSMS-G065-0002_Tiêu Chuẩn Treo Linh Kiện Cắt.pptx":
        "TSMS-G065-0002_절단 부품 걸이 표준_ko.pptx",
    "TSNC-G065-0001_a1_quy trình xử lý linh kiện không phù hợp.pptx":
        "TSNC-G065-0001_a1_부적합 부품 처리 절차_ko.pptx",
    "TSPO-G065-0001_A1_20240610_Quy trình hàn siêu âm.pptx":
        "TSPO-G065-0001_A1_20240610_초음파 용착 공정_ko.pptx",
    "TSWS-G065-0001_A1_20240610_Phương pháp quấn miếng bảo vệ linh kiện dập.pptx":
        "TSWS-G065-0001_A1_20240610_압착 부품 보호 테이프 감기 방법_ko.pptx",
    "TSWS-G065-0002_Tiêu chuẩn sử dụng thước đo điện tử (micromet)_a2.pptx":
        "TSWS-G065-0002_전자 측정기(마이크로미터) 사용 표준_a2_ko.pptx",
}

def is_vietnamese(text):
    return any(c in VN_CHARS for c in text)

def to_concise(text: str) -> str:
    t = text.strip()
    if len(t) < 3:
        return t
    for v in ACTION_VERBS:
        for ending in ['합니다', '한다', '하다']:
            t = re.sub(rf'([을를])\s*{v}{ending}\.?$', f' {v}.', t)
        t = re.sub(rf'([을를])\s*{v}해야\s*합니다\.?$', f' {v}해야 함.', t)
    endings = [
        (r'해야\s*합니다\.?$',   '해야 함.'),
        (r'해야\s*한다\.?$',     '해야 함.'),
        (r'하여야\s*합니다\.?$', '해야 함.'),
        (r'하십시오\.?$',        '할 것.'),
        (r'하세요\.?$',          '할 것.'),
        (r'해주세요\.?$',        '할 것.'),
        (r'합니다\.?$',          '함.'),
        (r'입니다\.?$',          '임.'),
        (r'됩니다\.?$',          '됨.'),
        (r'있습니다\.?$',        '있음.'),
        (r'없습니다\.?$',        '없음.'),
        (r'한다\.?$',            '함.'),
        (r'된다\.?$',            '됨.'),
        (r'있다\.?$',            '있음.'),
        (r'없다\.?$',            '없음.'),
        (r'이다\.?$',            '임.'),
    ]
    for pattern, replacement in endings:
        new_t = re.sub(pattern, replacement, t)
        if new_t != t:
            t = new_t
            break
    t = re.sub(r'([을를])\s*사용하여', ' 활용', t)
    t = re.sub(r'([을를])\s*이용하여', ' 활용', t)
    t = re.sub(r'\.{2,}$', '.', t)
    return t

def apply_term_rules(text: str) -> str:
    for vn, ko in TERM_RULES.items():
        text = text.replace(vn, ko)
    return text

def translate_text(text, translator):
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    try:
        result = translator.translate(text)
        time.sleep(0.05)
        if result:
            result = apply_term_rules(result)
            result = to_concise(result)
        return result if result else text
    except Exception:
        return text

def process_text_frame(tf, translator):
    for para in tf.paragraphs:
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator)
        elif len(para.runs) > 1:
            translated = translate_text(full_text, translator)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''

def process_shape(shape, translator):
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, translator)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, translator)
    if shape.shape_type == 6:
        for s in shape.shapes:
            process_shape(s, translator)

def translate_pptx(input_path, output_path):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, translator)
    prs.save(output_path)

def main():
    src_dir = Path(r"C:\hex_edit_project\VN_original\065_Tiêu chuẩn_Công đoạn_Hàn siêu âm")
    out_dir = Path(r"C:\베트남 티에스텍 업무\13. 절압팀\절차서 TC\065 표준_공정_초음파 용착 작업")
    out_dir.mkdir(parents=True, exist_ok=True)

    files = sorted([f for f in src_dir.iterdir() if f.suffix.lower() in ('.pptx', '.ppt')])
    total = len(files)
    print(f"총 {total}개 파일 번역 시작\n")

    errors = []
    for idx, f in enumerate(files, 1):
        out_name = FILENAME_MAP.get(f.name)
        if not out_name:
            print(f"[{idx}/{total}] 파일명 매핑 없음, 스킵: {f.name}")
            errors.append(f.name)
            continue

        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name}\n        → {out_name} 번역 중...")
        try:
            translate_pptx(str(f), str(out_path))
            print(f"[{idx}/{total}] 완료\n")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {e}\n")
            errors.append(f.name)

    print("=" * 60)
    print(f"번역 완료. 성공: {total - len(errors)}개 / 실패: {len(errors)}개")
    if errors:
        for e in errors:
            print(f"  - {e}")
    return errors

if __name__ == "__main__":
    main()
