import sys
import io
import re
import json
import time
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

VN_CHARS = set(
    'àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩ'
    'òóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđ'
    'ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨ'
    'ÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ'
)

TERM_RULES = {
    "Dập Tự Động": "압착기",
    "dập tự động": "압착기",
    "Dập tự động": "압착기",
}

ACTION_VERBS = [
    '확인', '수행', '준비', '진행', '관리', '적용', '설정', '실시',
    '사용', '기록', '교체', '점검', '조정', '검사', '실행', '처리',
    '보관', '표시', '부착', '제거', '측정', '삽입', '연결', '분리',
    '청소', '유지', '보전', '작성', '등록', '배치', '정리', '고정',
    '투입', '장착', '탈착', '조립', '해체', '세척', '건조', '포장',
]

FILENAME_MAP = {
    "TSCS-G080-0001_A1_20240610_CO NHIỆT_Phiếu kiểm tra ( sản xuất )_H-01.ppt":
        "TSCS-G080-0001_A1_20240610_열수축_검사표 (생산)_H-01_ko.pptx",
    "TSCS-G080-0002_A1_20240610_CO NHIỆT_Phiếu kiểm tra ( sản xuất )_H-02.ppt":
        "TSCS-G080-0002_A1_20240610_열수축_검사표 (생산)_H-02_ko.pptx",
    "TSCS-G080-0003_A1_20240610_Phiếu kiểm tra ( bảo trì ).ppt":
        "TSCS-G080-0003_A1_20240610_검사표 (보전)_ko.pptx",
    "TSIS-G080-0001_a2_20240822_Tiêu chuẩn tự kiểm tra khớp nối.pptx":
        "TSIS-G080-0001_a2_20240822_조인트 자가검사 표준_ko.pptx",
    "TSIS-G080-0002_a2_Tiêu chuẩn tự kiểm tra Ring terminal.pptx":
        "TSIS-G080-0002_a2_링 터미널 자가검사 표준_ko.pptx",
    "TSIS-G080-0003_Tiêu chuẩn tự kiểm tra của ống co nhiệt (sóng siêu âm)_a1.pptx":
        "TSIS-G080-0003_열수축 튜브 자가검사 표준 (초음파)_a1_ko.pptx",
    "TSMS_G080-0001_A0_20250217_Tiêu chuẩn quản lý ống co nhiệt_(+dán bàn co nhiệt).pptx":
        "TSMS_G080-0001_A0_20250217_열수축 튜브 관리 표준_(+게시판 부착)_ko.pptx",
    "TSMS-G080-0002_A1_20240610_Tiêu chuẩn nhận dạng phiếu sản xuất.pptx":
        "TSMS-G080-0002_A1_20240610_생산표 식별 표준_ko.pptx",
    "TSMS-G080-0003_A1_20240610_Tiêu chuẫn nhiệt độ và tốc độ máy co nhiệt HT-01+02(+dán bàn).pptx":
        "TSMS-G080-0003_A1_20240610_열수축기 HT-01+02 온도 및 속도 표준_(+게시판 부착)_ko.pptx",
    "TSMS-G080-0004_A1_20240610_Tiêu chuẫn nhiệt độ que đo nhiệt(DÁN BÀN ).pptx":
        "TSMS-G080-0004_A1_20240610_온도 프로브 온도 표준_(게시판 부착)_ko.pptx",
    "TSMS-G080-0005_Tiêu chuẩn phân biệt ống co nhiệt_a2.pptx":
        "TSMS-G080-0005_열수축 튜브 구분 표준_a2_ko.pptx",
    "TSMS-G080-0006_Tiêu chuẩn cốc bảo vệ_a3.pptx":
        "TSMS-G080-0006_보호 캡 표준_a3_ko.pptx",
    "TSMS-G080-0007_Tiêu chuẩn quản lý co nhiệt_a3.pptx":
        "TSMS-G080-0007_열수축 관리 표준_a3_ko.pptx",
    "TSMS-G080-0008_Tiêu chuẩn giữ vật phẩm cắt-dập (Khi cất giữ trong hộp)_a0.pptx":
        "TSMS-G080-0008_절단-압착 부품 보관 표준 (박스 보관 시)_a0_ko.pptx",
    "TSMS-G080-0009_A2_Tiêu Chuẩn Treo Linh Kiện Cắt – Dập.pptx":
        "TSMS-G080-0009_A2_절단-압착 부품 걸기 표준_ko.pptx",
    "TSMS-G080-0010_Tiêu chuẩn quản lý đặc tính đặc biệt (nhiệt độ  tốc độ)_a2.pptx":
        "TSMS-G080-0010_특수 특성 관리 표준 (온도 속도)_a2_ko.pptx",
    "TSNC-G080-0001_A1_20240610_Quy Trình Xử Lý Sản Phẩm Không Phù Hợp.pptx":
        "TSNC-G080-0001_A1_20240610_부적합 제품 처리 절차_ko.pptx",
    "TSPO-G080-0001_a3_Quy trình làm việc co nhiệt .pptx":
        "TSPO-G080-0001_a3_열수축 작업 절차_ko.pptx",
    "TSPO-G080-0002_Quy trình kiểm tra co nhiệt_a2.pptx":
        "TSPO-G080-0002_열수축 검사 절차_a2_ko.pptx",
    "TSWM-G080-0009_HT02_A3_20240622_Tiêu chuẩn co nhiệt (HT-02+01_TK1).pptx":
        "TSWM-G080-0009_HT02_A3_20240622_열수축 표준 (HT-02+01_TK1)_ko.pptx",
    "TSWM-G080-HT01_HT01_A5_20251110_Tiêu chuẩn co nhiệt (HT-01)(DÁN BÀN).pptx":
        "TSWM-G080-HT01_HT01_A5_20251110_열수축 표준 (HT-01)(게시판 부착)_ko.pptx",
    "TSWM-G080-HT02_HT02_A5_20251110_Tiêu chuẩn co nhiệt (HT-02)(DÁN BÀN).pptx":
        "TSWM-G080-HT02_HT02_A5_20251110_열수축 표준 (HT-02)(게시판 부착)_ko.pptx",
    "TSWS-G080_0001_A0_20240115_Phương pháp mở biểu đồ co nhiệt(DÁN BÀN).pptx":
        "TSWS-G080_0001_A0_20240115_열수축 차트 열기 방법_(게시판 부착)_ko.pptx",
    "TSWS-G080_0007_A0_20240115_Phương pháp scan vật tư co nhiệt.pptx":
        "TSWS-G080_0007_A0_20240115_열수축 자재 스캔 방법_ko.pptx",
    "TSWS-G080-0002_A0_20240610_Tiêu chuẩn sử dụng nhiệt kế.pptx":
        "TSWS-G080-0002_A0_20240610_온도계 사용 표준_ko.pptx",
    "TSWS-G080-0003_A1_20240610_Tiêu chuẩn quản lý của quá trình ép nhiệt (quản lý SPC)_làm lại theo TS.pptx":
        "TSWS-G080-0003_A1_20240610_열압착 공정 관리 표준 (SPC 관리)_ko.pptx",
    "TSWS-G080-0004_A1_20240610_Tiêu chuẩn chèn ống co nhiệt.pptx":
        "TSWS-G080-0004_A1_20240610_열수축 튜브 삽입 표준_ko.pptx",
    "TSWS-G080-0005_A2_20251110_Tiêu chuẩn mẫu thử ép nhiệt.pptx":
        "TSWS-G080-0005_A2_20251110_열압착 시험 샘플 표준_ko.pptx",
    "TSWS-G080-0006_A1_20240610_Tiêu Chuẩn Khắc Phục (Khi Máy Xảy Ra Vấn Đề).pptx":
        "TSWS-G080-0006_A1_20240610_복구 표준 (설비 문제 발생 시)_ko.pptx",
    "TSWS-G080-0008_A0_Tiêu chuẩn sử dụng nhiệt kế _máy khò nhiệt.pptx":
        "TSWS-G080-0008_A0_온도계 사용 표준_열풍기_ko.pptx",
    "TSWS-G080-RG3 CONT THETA3-00001_a1_TIÊU CHUẨN LÀM VIỆC SÚNG BẮN NHIỆT.pptx":
        "TSWS-G080-RG3_CONT_THETA3-00001_a1_열풍건 작업 표준_ko.pptx",
}

def load_translation_dict(dict_path):
    try:
        with open(dict_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {}

def apply_dict(text, trans_dict):
    for vn, ko in trans_dict.items():
        if vn in text:
            text = text.replace(vn, ko)
    return text

def is_vietnamese(text):
    return any(c in VN_CHARS for c in text)

def to_concise(text: str) -> str:
    t = text.strip()
    if len(t) < 3:
        return t
    for v in ACTION_VERBS:
        for ending in ['합니다', '한다', '하다']:
            t = re.sub(rf'([을를])\s*{v}{ending}\.?$', f' {v}.', t)
        t = re.sub(rf'([을를])\s*{v}해야\s*합니다\.?$', f' {v}해야 함.', t)
    endings = [
        (r'해야\s*합니다\.?$',   '해야 함.'),
        (r'해야\s*한다\.?$',     '해야 함.'),
        (r'하여야\s*합니다\.?$', '해야 함.'),
        (r'하십시오\.?$',        '할 것.'),
        (r'하세요\.?$',          '할 것.'),
        (r'해주세요\.?$',        '할 것.'),
        (r'합니다\.?$',          '함.'),
        (r'입니다\.?$',          '임.'),
        (r'됩니다\.?$',          '됨.'),
        (r'있습니다\.?$',        '있음.'),
        (r'없습니다\.?$',        '없음.'),
        (r'한다\.?$',            '함.'),
        (r'된다\.?$',            '됨.'),
        (r'있다\.?$',            '있음.'),
        (r'없다\.?$',            '없음.'),
        (r'이다\.?$',            '임.'),
    ]
    for pattern, replacement in endings:
        new_t = re.sub(pattern, replacement, t)
        if new_t != t:
            t = new_t
            break
    t = re.sub(r'([을를])\s*사용하여', ' 활용', t)
    t = re.sub(r'([을를])\s*이용하여', ' 활용', t)
    t = re.sub(r'\.{2,}$', '.', t)
    return t

def apply_term_rules(text: str) -> str:
    for vn, ko in TERM_RULES.items():
        text = text.replace(vn, ko)
    return text

def translate_text(text, translator, trans_dict):
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    applied = apply_dict(text, trans_dict)
    if not is_vietnamese(applied):
        return apply_term_rules(to_concise(applied))
    try:
        result = translator.translate(applied)
        time.sleep(0.05)
        if result:
            result = apply_term_rules(result)
            result = to_concise(result)
        return result if result else text
    except Exception:
        return text

def process_text_frame(tf, translator, trans_dict):
    for para in tf.paragraphs:
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator, trans_dict)
        elif len(para.runs) > 1:
            translated = translate_text(full_text, translator, trans_dict)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''

def process_shape(shape, translator, trans_dict):
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, translator, trans_dict)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, translator, trans_dict)
    if shape.shape_type == 6:
        for s in shape.shapes:
            process_shape(s, translator, trans_dict)

def translate_pptx(input_path, output_path, trans_dict):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, translator, trans_dict)
    prs.save(output_path)

def main():
    src_dir = Path(r"C:\hex_edit_project\VN_original\080_Tiêu chuẩn_Công đoạn_Co nhiệt\080_Tiêu chuẩn_Công đoạn_Co nhiệt")
    out_dir = Path(r"C:\hex_edit_project\KO_result")
    out_dir.mkdir(parents=True, exist_ok=True)

    dict_path = Path(r"C:\hex_edit_project\translation_dict.json")
    trans_dict = load_translation_dict(dict_path)
    print(f"번역 사전 {len(trans_dict)}개 항목 로드 완료\n")

    files = sorted([f for f in src_dir.iterdir() if f.suffix.lower() in ('.pptx', '.ppt') and f.name != 'Thumbs.db'])
    total = len(files)
    print(f"총 {total}개 파일 번역 시작\n")

    errors = []
    skipped = []
    for idx, f in enumerate(files, 1):
        out_name = FILENAME_MAP.get(f.name)
        if not out_name:
            print(f"[{idx}/{total}] 파일명 매핑 없음, 스킵: {f.name}")
            skipped.append(f.name)
            continue

        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name}\n        → {out_name} 번역 중...")
        try:
            translate_pptx(str(f), str(out_path), trans_dict)
            print(f"[{idx}/{total}] 완료\n")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {e}\n")
            errors.append(f.name)

    print("=" * 60)
    success = total - len(errors) - len(skipped)
    print(f"번역 완료. 성공: {success}개 / 오류: {len(errors)}개 / 스킵: {len(skipped)}개")
    if errors:
        print("오류 파일:")
        for e in errors:
            print(f"  - {e}")
    if skipped:
        print("스킵 파일:")
        for s in skipped:
            print(f"  - {s}")
    return errors, skipped

if __name__ == "__main__":
    main()
