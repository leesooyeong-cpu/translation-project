import os
import sys
import io
import time
import re
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

VN_CHARS = set('àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩòóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ')

TERM_MAP = {
    'Dập Tự Động': '압착기',
    'dập tự động': '압착기',
    'DẬP TỰ ĐỘNG': '압착기',
}

def apply_terms(text):
    for vn, ko in TERM_MAP.items():
        text = text.replace(vn, ko)
    return text

def is_vietnamese(text):
    return any(c in VN_CHARS for c in text)

def translate_text(text, translator):
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    try:
        result = translator.translate(text)
        time.sleep(0.05)
        translated = result if result else text
        return apply_terms(translated)
    except Exception:
        return text

def process_text_frame(tf, translator):
    for para in tf.paragraphs:
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator)
        elif len(para.runs) > 1:
            translated = translate_text(full_text, translator)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''

def translate_pptx(input_path, output_path):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                process_text_frame(shape.text_frame, translator)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_text_frame(cell.text_frame, translator)
            if shape.shape_type == 6:
                for s in shape.shapes:
                    if s.has_text_frame:
                        process_text_frame(s.text_frame, translator)

    prs.save(output_path)

def main():
    src_dir = Path(r"C:\hex_edit_project\VN_original\042_Tiêu chuẩn máy dập bán tự động 30TONS")
    out_dir = Path(r"C:\hex_edit_project\KO_result")
    out_dir.mkdir(exist_ok=True)

    pptx_files = sorted([f for f in src_dir.glob("*.pptx")])
    ppt_files = sorted([f for f in src_dir.glob("*.ppt")])

    total = len(pptx_files)
    skipped = [f.name for f in ppt_files]

    print(f"총 {total}개 pptx 파일 번역 시작 (ppt 제외: {len(ppt_files)}개)\n")

    errors = []
    for idx, f in enumerate(pptx_files, 1):
        out_name = f.stem + "_ko" + f.suffix
        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name} → 번역 중...")
        try:
            translate_pptx(str(f), str(out_path))
            print(f"[{idx}/{total}] 완료: {out_name}")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {f.name} → {e}")
            errors.append(f.name)

    print("\n=== 번역 완료 ===")
    print(f"성공: {total - len(errors)}개")
    if errors:
        print(f"오류: {errors}")
    if skipped:
        print(f"번역 불가 (.ppt 구형 포맷): {skipped}")

if __name__ == "__main__":
    main()
