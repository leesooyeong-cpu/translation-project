import os
import sys
import io
import time
from pathlib import Path
from pptx import Presentation
from deep_translator import GoogleTranslator

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

VN_CHARS = set('àáảãạăắặằẳẵâấầẩẫậèéẹẻẽêếềệểễìíịỉĩòóọỏõôốồổỗộơớờởỡợùúụủũưứừửữựỳýỵỷỹđÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬÈÉẸẺẼÊẾỀỆỂỄÌÍỊỈĨÒÓỌỎÕÔỐỒỔỖỘƠỚỜỞỠỢÙÚỤỦŨƯỨỪỬỮỰỲÝỴỶỸĐ')

def is_vietnamese(text):
    return any(c in VN_CHARS for c in text)

def translate_text(text, translator):
    if not text or not text.strip() or not is_vietnamese(text):
        return text
    try:
        result = translator.translate(text)
        time.sleep(0.05)
        return result if result else text
    except Exception:
        return text

def process_text_frame(tf, translator):
    for para in tf.paragraphs:
        # 단락 전체 텍스트 수집
        full_text = ''.join(r.text for r in para.runs)
        if not is_vietnamese(full_text):
            continue
        # run이 1개면 바로 교체
        if len(para.runs) == 1:
            para.runs[0].text = translate_text(para.runs[0].text, translator)
        elif len(para.runs) > 1:
            # 전체 번역 후 첫 run에 넣고 나머지 비움
            translated = translate_text(full_text, translator)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ''

def translate_pptx(input_path, output_path):
    translator = GoogleTranslator(source='vi', target='ko')
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                process_text_frame(shape.text_frame, translator)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_text_frame(cell.text_frame, translator)
            # GroupShape (flowchart 구성 도형 포함)
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for s in shape.shapes:
                    if s.has_text_frame:
                        process_text_frame(s.text_frame, translator)

    prs.save(output_path)

def main():
    src_dir = Path(r"C:\hex_edit_project\VN_original\020_Tiêu chuẩn_Công đoạn_Cắt dập tự động")
    out_dir = Path(r"C:\hex_edit_project\KO_result")
    out_dir.mkdir(exist_ok=True)

    files = sorted([f for f in src_dir.glob("*.pptx") if not f.name.startswith("(한글)")])
    total = len(files)
    print(f"총 {total}개 파일 번역 시작\n")

    for idx, f in enumerate(files, 1):
        out_name = f.stem + "_ko" + f.suffix
        out_path = out_dir / out_name
        print(f"[{idx}/{total}] {f.name} → 번역 중...")
        try:
            translate_pptx(str(f), str(out_path))
            print(f"[{idx}/{total}] 완료: {out_name}")
        except Exception as e:
            print(f"[{idx}/{total}] 오류: {f.name} → {e}")

    print("\n전체 번역 완료.")

if __name__ == "__main__":
    main()
